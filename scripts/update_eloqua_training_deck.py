#!/usr/bin/env python3
"""Add speaker notes and populate slides 4 & 5 in the Eloqua training deck."""

from __future__ import annotations

import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

WORKSPACE = Path(__file__).resolve().parents[1]
DECK = WORKSPACE / "Eloqua Lead Scoring Training.pptx"
XLSX = WORKSPACE / "FY26Q4-July2026-Weight Changes.xlsx"

HEADER_FILL = RGBColor(0x00, 0x52, 0x8B)
ALT_FILL = RGBColor(0xF2, 0xF4, 0xF6)
TEXT_COLOR = RGBColor(0x58, 0x59, 0x5B)

SPEAKER_NOTES = {
    1: (
        "Welcome the audience and introduce yourself.\n"
        "Set expectations: this session explains how Eloqua scores leads, why the model was recalibrated, "
        "and what it means for the leads sales and VDC receive.\n"
        "Timing: ~30 minutes for sales audiences; ~55 minutes if covering the full technical flow."
    ),
    2: (
        "Read both objectives aloud, then pause for engagement.\n"
        "Ask: 'How many of you have chased a lead that went nowhere?' — this builds buy-in.\n"
        "Emphasise: we are reducing passive-interest leads and rebuilding sales trust with higher-intent delivery."
    ),
    3: (
        "Explain the two-score model before showing any tables.\n"
        "Profile (A–D) = explicit fit data — who the prospect is (title, department, email type).\n"
        "Engagement (1–4) = implicit intent data — what they did and how recently.\n"
        "Key message: both scores work together. A strong profile with low engagement is a nurture candidate; "
        "high engagement with weak profile still warrants review."
    ),
    4: (
        "Walk through the four profile categories — each weighted at 25%.\n"
        "Point to the grade thresholds: A (75–100) = excellent ICP fit; B (50–74) = good fit; "
        "C (25–49) = marginal; D (0–24) = poor fit.\n"
        "Do not read every department row — highlight that MIS/IT, Network Management, and Operations score "
        "17.5 pts (70%), while most other departments score 6.25 pts (25%).\n"
        "Corporate email = 25 pts; public email = 1.25 pts — this is a strong differentiator."
    ),
    5: (
        "Introduce the 12 engagement activity categories and their maximum weights from the FY26 Q4 model.\n"
        "Highlight the highest-weight activities: Demo Scheduled (25%) and Event Attendance (25%).\n"
        "Explain recency decay: scores are highest within 7 days, lower at 14 days, lowest at 30 days.\n"
        "Explain frequency: 'More than 1 time' scores higher than 'At least 1 time' within the same window.\n"
        "Engagement grades: 4 (76–100) = hot; 3 (51–75) = warm-high; 2 (26–50) = warm; 1 (0–25) = cold."
    ),
    6: (
        "Do not read every bullet — group into four buckets:\n"
        "1) Events (Cvent, RainFocus, booth scans)\n"
        "2) Forms (demo, callback, webinar registration)\n"
        "3) Integrations (paid media, video, inbound calls/chats)\n"
        "4) Digital (email clicks, website, landing pages)\n"
        "Key message: not all activity is equal — a demo scores far higher than a landing page visit."
    ),
    7: (
        "Orient the audience to the end-to-end data flow before the detailed processing slide.\n"
        "Say: 'Activity enters Eloqua, is scored, validated, and routed to VDC when thresholds are met.'\n"
        "For sales audiences, spend 1 minute here then skip ahead to slide 9."
    ),
    8: (
        "Technical audience slide — walk through left to right.\n"
        "Key steps: (1) transactions prepared with enhanced detail, (2) scoring model runs, "
        "(3) data written to Scored Leads CDO, (4) segment created, (5) Program 389 determines next steps, "
        "(6) high-score contacts sent to VDC with last-touch data.\n"
        "Pause at decision points: High Score (Yes/No), CCID present, Valid Lead, AU Public Sector check.\n"
        "For sales audiences, summarise in one sentence and move on."
    ),
    9: (
        "Tell Tony's story before revealing the table.\n"
        "Tony is an IT Manager at Intuitech Technology Ltd, exploring AI for his company.\n"
        "Build the score live:\n"
        "  Job Level: Manager → 6.25\n"
        "  Decision Maker: Technical → 18.75\n"
        "  Department: MIS/IT → 17.5\n"
        "  Email: Corporate → 25.0\n"
        "  Total = 67.5 → Profile Grade B (good ICP fit).\n"
        "Note: title says 'Director' but narrative describes a Manager — clarify if asked."
    ),
    10: (
        "Walk through Tony's engagement journey chronologically:\n"
        "1) Email click 10 days ago with no action → low score due to time decay.\n"
        "2) Returned 6 days later, watched 50% of video → 9 pts.\n"
        "3) Today: attended webinar (13.5 pts) and scheduled a demo (25 pts).\n"
        "Total engagement ≈ 53 → Grade 2 (warm — building intent).\n"
        "End with the punchline: 'After watching the video, Tony scheduled a demo — that's the lead we want sales to call.'\n"
        "Open floor for questions."
    ),
}


def read_xlsx_rows(path: Path) -> tuple[list[list[str]], list[list[str]]]:
    with zipfile.ZipFile(path) as zf:
        shared: list[str] = []
        root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
        ns = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        for item in root.findall(".//m:si", ns):
            texts = [node.text or "" for node in item.findall(".//m:t", ns)]
            shared.append("".join(texts))

        def sheet_rows(sheet_name: str) -> list[list[str]]:
            xml = ET.fromstring(zf.read(f"xl/worksheets/{sheet_name}.xml"))
            rows: list[list[str]] = []
            for row in xml.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}row"):
                values: list[str] = []
                for cell in row.findall("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}c"):
                    cell_type = cell.get("t")
                    value = cell.find("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}v")
                    if value is None:
                        values.append("")
                    elif cell_type == "s":
                        values.append(shared[int(value.text)])
                    else:
                        values.append(value.text or "")
                rows.append(values)
            return rows

        return sheet_rows("sheet1"), sheet_rows("sheet2")


def parse_profile(rows: list[list[str]]) -> list[dict]:
    categories: list[dict] = []
    current: dict | None = None
    for row in rows[1:]:
        if not row:
            continue
        if len(row) >= 5 and row[3]:
            current = {
                "name": row[0],
                "weight": row[3],
                "values": [{"value": row[1], "score_pct": row[2], "points": row[4]}],
            }
            categories.append(current)
        elif current and row[0]:
            score_pct = row[1] if len(row) > 1 else ""
            points = row[2] if len(row) > 2 else ""
            current["values"].append({"value": row[0], "score_pct": score_pct, "points": points})
    return categories


def parse_engagement(rows: list[list[str]]) -> list[dict]:
    categories: list[dict] = []
    current_name = ""
    for row in rows[1:]:
        if not row:
            continue
        if row[0]:
            current_name = row[0]
        if current_name and row[5] and not any(c["name"] == current_name for c in categories):
            categories.append({"name": current_name, "weight": row[5]})
    return categories


def set_font(paragraph, *, size: int = 10, bold: bool = False, color: RGBColor | None = None):
    font = paragraph.font
    font.name = "Segoe UI"
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = color


def style_table_cell(cell, *, header: bool = False, alt: bool = False):
    if header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        text_color = RGBColor(0xFF, 0xFF, 0xFF)
        bold = True
        size = 10
    else:
        if alt:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ALT_FILL
        text_color = TEXT_COLOR
        bold = False
        size = 9
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER if header else PP_ALIGN.LEFT
        set_font(paragraph, size=size, bold=bold, color=text_color)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for idx, width_val in enumerate(col_widths):
            table.columns[idx].width = width_val

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        style_table_cell(cell, header=True)

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            style_table_cell(cell, alt=row_idx % 2 == 0)
    return table_shape


def add_textbox(slide, left, top, width, height, text, size=11, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    set_font(p, size=size, bold=bold, color=TEXT_COLOR)
    return box


def build_profile_rows(categories: list[dict]) -> list[list[str]]:
    rows: list[list[str]] = []
    for category in categories:
        first = True
        for item in category["values"]:
            if first:
                rows.append(
                    [
                        category["name"],
                        item["value"],
                        f"{item['score_pct']}%",
                        f"{category['weight']}%",
                        item["points"],
                    ]
                )
                first = False
            else:
                rows.append(["", item["value"], f"{item['score_pct']}%", "", item["points"]])
    return rows


def populate_slide_4(slide, categories: list[dict]):
    grade_rows = [
        ["A", "75 – 100", "Excellent ICP fit"],
        ["B", "50 – 74", "Good fit — worth pursuing"],
        ["C", "25 – 49", "Marginal fit — nurture"],
        ["D", "0 – 24", "Poor fit"],
    ]
    add_table(
        slide,
        Inches(8.9),
        Inches(1.0),
        Inches(4.0),
        Inches(1.55),
        ["Grade", "Score Range", "Meaning"],
        grade_rows,
        col_widths=[Inches(0.55), Inches(1.1), Inches(2.35)],
    )

    add_textbox(
        slide,
        Inches(8.9),
        Inches(2.65),
        Inches(4.0),
        Inches(0.55),
        "Four categories, 25% weight each. Maximum profile score = 100.",
        size=10,
        bold=True,
    )

    profile_rows = build_profile_rows(categories)
    add_table(
        slide,
        Inches(0.35),
        Inches(1.0),
        Inches(8.35),
        Inches(6.2),
        ["Category", "Field Value", "Score %", "Weight", "Points"],
        profile_rows,
        col_widths=[Inches(1.55), Inches(2.55), Inches(0.85), Inches(0.75), Inches(0.85)],
    )


def populate_slide_5(slide, categories: list[dict]):
    grade_rows = [
        ["4", "76 – 100", "Hot — immediate follow-up"],
        ["3", "51 – 75", "Warm-high — prioritise"],
        ["2", "26 – 50", "Warm — nurture and monitor"],
        ["1", "0 – 25", "Cold — low-touch"],
    ]
    add_table(
        slide,
        Inches(8.9),
        Inches(1.0),
        Inches(4.0),
        Inches(1.55),
        ["Grade", "Score Range", "Meaning"],
        grade_rows,
        col_widths=[Inches(0.55), Inches(1.1), Inches(2.35)],
    )

    add_textbox(
        slide,
        Inches(8.9),
        Inches(2.65),
        Inches(4.0),
        Inches(1.2),
        "Recency windows: 7, 14, and 30 days.\n"
        "Frequency: 'More than 1 time' scores higher than 'At least 1 time' within the same window.\n"
        "Maximum engagement score = 100.",
        size=10,
    )

    engagement_rows = [
        [cat["name"], f"{cat['weight']}%", "7 / 14 / 30 days", "At least 1 / More than 1"]
        for cat in categories
    ]
    add_table(
        slide,
        Inches(0.35),
        Inches(1.0),
        Inches(8.35),
        Inches(6.2),
        ["Activity Category", "Max Weight", "Recency Windows", "Frequency Rules"],
        engagement_rows,
        col_widths=[Inches(2.5), Inches(1.1), Inches(1.5), Inches(2.0)],
    )


def set_speaker_notes(slide, notes: str):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.clear()
    text_frame.text = notes


def update_deck():
    engagement_rows, profile_rows = read_xlsx_rows(XLSX)
    profile_categories = parse_profile(profile_rows)
    engagement_categories = parse_engagement(engagement_rows)

    prs = Presentation(DECK)

    populate_slide_4(prs.slides[3], profile_categories)
    populate_slide_5(prs.slides[4], engagement_categories)

    for slide_num, notes in SPEAKER_NOTES.items():
        set_speaker_notes(prs.slides[slide_num - 1], notes)

    prs.save(DECK)
    print(f"Updated {DECK}")


if __name__ == "__main__":
    update_deck()

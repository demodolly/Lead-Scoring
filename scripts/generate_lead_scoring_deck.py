#!/usr/bin/env python3
"""Generate Lead Scoring training PowerPoint from FY26Q4 weight model."""

from __future__ import annotations

import zipfile
import xml.etree.ElementTree as ET
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

WORKSPACE = Path(__file__).resolve().parents[1]
XLSX = WORKSPACE / "FY26Q4-July2026-Weight Changes.xlsx"
OUTPUT = WORKSPACE / "Lead-Scoring-Display-Guide.pptx"

# Cisco-inspired palette
CISCO_BLUE = RGBColor(0x00, 0x7B, 0xC7)
CISCO_DARK = RGBColor(0x00, 0x52, 0x8B)
CISCO_TEAL = RGBColor(0x00, 0x96, 0xA6)
CISCO_GREEN = RGBColor(0x6B, 0xB3, 0x45)
CISCO_AMBER = RGBColor(0xF5, 0xA6, 0x23)
CISCO_GRAY = RGBColor(0x58, 0x59, 0x5B)
CISCO_LIGHT = RGBColor(0xF2, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def read_xlsx_rows(path: Path) -> tuple[list[list[str]], list[list[str]]]:
    with zipfile.ZipFile(path) as zf:
        shared: list[str] = []
        root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
        ns = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
        for item in root.findall(".//m:si", ns):
            texts = [node.text or "" for node in item.findall(".//m:t", ns)]
            shared.append("".join(texts))

        def sheet_rows(sheet_name: str) -> list[list[str]]:
            xml = ET.fromstring(zf.read(f"xl/worksheets/{sheet_name}.xml"))
            rows: list[list[str]] = []
            for row in xml.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}row"):
                values: list[str] = []
                for cell in row.findall("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}c"):
                    cell_type = cell.get("t")
                    value = cell.find("{http://schemas.openxmlformats.org/spreadsheetml/2006/main}v")
                    if value is None:
                        values.append("")
                    elif cell_type == "s":
                        values.append(shared[int(value.text)])
                    else:
                        values.append(value.text or "")
                rows.append(values)
            return rows

        return sheet_rows("sheet1"), sheet_rows("sheet2")


def parse_engagement(rows: list[list[str]]) -> list[dict]:
    categories: dict[str, dict] = {}
    current = ""
    for row in rows[1:]:
        if not row:
            continue
        if row[0]:
            current = row[0]
        if current not in categories:
            categories[current] = {"name": current, "weight": row[5], "rules": []}
        categories[current]["rules"].append(
            {
                "times": row[2],
                "time_frame": row[3],
                "score_pct": row[4],
                "final_score": row[6] if len(row) > 6 else "",
            }
        )
    return list(categories.values())


def parse_profile(rows: list[list[str]]) -> list[dict]:
    categories: list[dict] = []
    current: dict | None = None
    for row in rows[1:]:
        if not row:
            continue
        if len(row) >= 5 and row[3]:
            current = {
                "name": row[0],
                "weight": row[3],
                "values": [{"value": row[1], "score_pct": row[2], "points": row[4]}],
            }
            categories.append(current)
        elif current and row[0]:
            score_pct = row[1] if len(row) > 1 else ""
            points = row[2] if len(row) > 2 else ""
            current["values"].append({"value": row[0], "score_pct": score_pct, "points": points})
    return categories


def set_run(font, *, size: int, bold: bool = False, color: RGBColor | None = None):
    font.name = "Segoe UI"
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = color


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CISCO_DARK
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    set_run(p.font, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(12.4), Inches(0.3))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        set_run(sp.font, size=13, color=WHITE)


def add_bullets(slide, left, top, width, height, items, font_size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.font, size=font_size, color=CISCO_GRAY)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = CISCO_DARK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            set_run(p.font, size=11, bold=True, color=WHITE)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CISCO_LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                set_run(p.font, size=10, color=CISCO_GRAY)
    return table


def add_score_pill(slide, left, top, width, height, label, score, max_score, accent):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.text = label
    p1.alignment = PP_ALIGN.CENTER
    set_run(p1.font, size=14, bold=True, color=WHITE)

    p2 = tf.add_paragraph()
    p2.text = f"{score} / {max_score}"
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2.font, size=24, bold=True, color=WHITE)


def add_progress_bar(slide, left, top, width, fraction, color):
    track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.22))
    track.fill.solid()
    track.fill.fore_color.rgb = RGBColor(0xE5, 0xE8, 0xEB)
    track.line.fill.background()

    fill_width = max(Inches(0.1), int(width * max(0.0, min(1.0, fraction))))
    fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, fill_width, Inches(0.22))
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()


def build_deck() -> Presentation:
    engagement_rows, profile_rows = read_xlsx_rows(XLSX)
    engagement = parse_engagement(engagement_rows)
    profile = parse_profile(profile_rows)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Title
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CISCO_DARK
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Lead Scoring Display Guide"
    set_run(tp.font, size=40, bold=True, color=WHITE)

    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.8))
    sp = subtitle.text_frame.paragraphs[0]
    sp.text = "How to read Engagement and Profile scores  |  FY26 Q4 (July 2026)"
    set_run(sp.font, size=18, color=WHITE)

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Lead Scoring Training  •  Cisco Marketing"
    set_run(fp.font, size=12, color=WHITE)

    # Slide 2: Overview
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "How Lead Scoring Works", "Two complementary scores combine into one actionable lead priority")

    add_score_pill(slide, Inches(1.0), Inches(1.55), Inches(3.2), Inches(1.1), "Engagement", 45, 100, CISCO_BLUE)
    add_score_pill(slide, Inches(5.05), Inches(1.55), Inches(3.2), Inches(1.1), "Profile", 27, 100, CISCO_TEAL)
    add_score_pill(slide, Inches(9.1), Inches(1.55), Inches(3.2), Inches(1.1), "Lead Score", 72, 100, CISCO_GREEN)

    plus = slide.shapes.add_textbox(Inches(4.25), Inches(1.85), Inches(0.5), Inches(0.4))
    plus.text_frame.paragraphs[0].text = "+"
    set_run(plus.text_frame.paragraphs[0].font, size=28, bold=True, color=CISCO_GRAY)

    eq = slide.shapes.add_textbox(Inches(8.3), Inches(1.85), Inches(0.5), Inches(0.4))
    eq.text_frame.paragraphs[0].text = "="
    set_run(eq.text_frame.paragraphs[0].font, size=28, bold=True, color=CISCO_GRAY)

    cards = [
        ("Engagement Score", "Measures intent", "What did the lead do, and how recently?", "12 behavioral signals with time decay"),
        ("Profile Score", "Measures fit", "Who is the lead based on firmographics?", "4 ICP attributes at 25% each"),
    ]
    for idx, (title_text, tag, question, detail) in enumerate(cards):
        left = Inches(0.9 + idx * 6.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.0), Inches(5.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CISCO_LIGHT
        card.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE)

        box = slide.shapes.add_textbox(left + Inches(0.25), Inches(3.2), Inches(5.3), Inches(3.1))
        tf = box.text_frame
        tf.clear()
        lines = [
            (title_text, 20, True, CISCO_DARK),
            (tag, 14, True, CISCO_BLUE if idx == 0 else CISCO_TEAL),
            ("", 6, False, CISCO_GRAY),
            (question, 15, False, CISCO_GRAY),
            (detail, 13, False, CISCO_GRAY),
        ]
        for i, (text, size, bold, color) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.space_after = Pt(6)
            set_run(p.font, size=size, bold=bold, color=color)

    # Slide 3: Sample score card
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Sample Lead Score Card", "What sales and marketing should see at a glance")

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.7))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE)

    headline = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(7.5), Inches(0.7))
    hp = headline.text_frame.paragraphs[0]
    hp.text = "Lead Score: 72 / 100"
    set_run(hp.font, size=30, bold=True, color=CISCO_DARK)

    tier = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.65), Inches(2.8), Inches(0.55))
    tier.fill.solid()
    tier.fill.fore_color.rgb = CISCO_GREEN
    tier.line.fill.background()
    tfp = tier.text_frame.paragraphs[0]
    tfp.text = "Hot Lead"
    tfp.alignment = PP_ALIGN.CENTER
    set_run(tfp.font, size=14, bold=True, color=WHITE)

    add_score_pill(slide, Inches(1.0), Inches(2.45), Inches(3.5), Inches(0.95), "Engagement", 45, 100, CISCO_BLUE)
    add_score_pill(slide, Inches(4.85), Inches(2.45), Inches(3.5), Inches(0.95), "Profile", 27, 100, CISCO_TEAL)
    add_progress_bar(slide, Inches(1.0), Inches(3.55), Inches(7.35), 0.45, CISCO_BLUE)
    add_progress_bar(slide, Inches(1.0), Inches(4.05), Inches(7.35), 0.27, CISCO_TEAL)

    add_bullets(
        slide,
        Inches(1.0),
        Inches(4.55),
        Inches(11.2),
        Inches(2.2),
        [
            "Top drivers:",
            "• Demo scheduled (last 7 days) — 25 / 25 pts",
            "• Decision Maker: TDM — 25 / 25 pts",
            "• Event attendance (last 14 days) — 17.5 / 25 pts",
            "Tip: Show points earned vs. category maximum so reps understand both strength and headroom.",
        ],
        font_size=15,
    )

    # Slide 4: Engagement weights
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Engagement Score Breakdown", "Category weights from FY26 Q4 model (max 100 points total)")

    eng_table_rows = [[cat["name"], f"{cat['weight']}%", "7 / 14 / 30 days", "At least 1 / More than 1"] for cat in engagement]
    add_table(
        slide,
        Inches(0.55),
        Inches(1.35),
        Inches(12.2),
        Inches(5.6),
        ["Activity Category", "Max Weight", "Recency Windows", "Frequency Rules"],
        eng_table_rows,
        col_widths=[Inches(3.4), Inches(1.5), Inches(2.2), Inches(2.4)],
    )

    note = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(12.0), Inches(0.4))
    np = note.text_frame.paragraphs[0]
    np.text = "Display tip: Use a horizontal bar per category showing points earned vs. max (e.g., Demo 25/25)."
    set_run(np.font, size=11, color=CISCO_GRAY)

    # Slide 5: Engagement decay example
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Engagement Recency Decay", "Example: Demo (Scheduled) — 25% category weight")

    demo_rules = next(cat for cat in engagement if "Demo" in cat["name"])["rules"][:6]
    decay_rows = [[rule["times"], rule["time_frame"], f"{rule['score_pct']}%", rule["final_score"]] for rule in demo_rules]
    add_table(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(5.8),
        Inches(3.2),
        ["Frequency", "Time Frame", "Score %", "Points"],
        decay_rows,
        col_widths=[Inches(1.8), Inches(1.2), Inches(1.0), Inches(1.0)],
    )

    legend = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5))
    ltf = legend.text_frame
    ltf.clear()
    legend_lines = [
        ("How to visualize decay", 18, True, CISCO_DARK),
        ("", 8, False, CISCO_GRAY),
        ("7 days  = highest points (green)", 14, False, CISCO_GREEN),
        ("14 days = medium points (amber)", 14, False, CISCO_AMBER),
        ("30 days = lower points (gray)", 14, False, CISCO_GRAY),
        ("", 8, False, CISCO_GRAY),
        ("More than 1 time scores higher than at least 1 time within the same window.", 13, False, CISCO_GRAY),
        ("Show an expiry badge on active signals (e.g., 'expires in 5 days').", 13, False, CISCO_GRAY),
    ]
    for i, (text, size, bold, color) in enumerate(legend_lines):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = text
        p.space_after = Pt(6)
        set_run(p.font, size=size, bold=bold, color=color)

  # Slide 6: Profile categories
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Profile Score Breakdown", "ICP fit across four equal-weight categories (25% each)")

    profile_summary = []
    for cat in profile:
        top_values = sorted(cat["values"], key=lambda v: float(v["points"] or 0), reverse=True)[:3]
        examples = ", ".join(f"{v['value']} ({v['points']})" for v in top_values)
        profile_summary.append([cat["name"], f"{cat['weight']}%", examples])

    add_table(
        slide,
        Inches(0.55),
        Inches(1.45),
        Inches(12.2),
        Inches(2.2),
        ["Profile Category", "Weight", "Top-scoring examples (points)"],
        profile_summary,
        col_widths=[Inches(2.5), Inches(1.2), Inches(7.8)],
    )

    checklist = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.0), Inches(5.8), Inches(2.8))
    checklist.fill.solid()
    checklist.fill.fore_color.rgb = CISCO_LIGHT
    checklist.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE)

    add_bullets(
        slide,
        Inches(0.95),
        Inches(4.2),
        Inches(5.3),
        Inches(2.4),
        [
            "Sample profile fit (27/100):",
            "✓ Decision Maker: TDM — 25/25",
            "✓ Job Level: Director — 18.75/25",
            "○ Job Department: Engineering — 6.25/25",
            "✓ Email Type: Corporate — 25/25",
        ],
        font_size=14,
    )

    add_bullets(
        slide,
        Inches(6.9),
        Inches(4.2),
        Inches(5.5),
        Inches(2.4),
        [
            "Display tips:",
            "• Use a checklist, not a dense lookup table",
            "• Green = matched and contributing",
            "• Gray = missing or low-value attribute",
            "• Helps reps answer: 'Is this person worth pursuing?'",
        ],
        font_size=14,
    )

    # Slide 7: Tier guidance
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Recommended Score Tiers", "Suggested thresholds for prioritization (adjust per campaign)")

    add_table(
        slide,
        Inches(1.2),
        Inches(1.8),
        Inches(10.8),
        Inches(2.0),
        ["Tier", "Lead Score Range", "Recommended Action"],
        [
            ["Hot", "70 – 100", "Immediate sales follow-up; assign to rep"],
            ["Warm", "40 – 69", "Nurture with targeted content; monitor engagement"],
            ["Cold", "0 – 39", "Keep in marketing automation; low-touch outreach"],
        ],
        col_widths=[Inches(1.5), Inches(2.2), Inches(6.5)],
    )

    add_bullets(
        slide,
        Inches(1.2),
        Inches(4.2),
        Inches(10.8),
        Inches(2.5),
        [
            "Key takeaways for any score display:",
            "1. Show composite score + Engagement + Profile side by side",
            "2. Always include top 3 drivers explaining why the score changed",
            "3. Label recency clearly (7d / 14d / 30d) for engagement signals",
            "4. Show points earned vs. maximum per category",
            "5. Use color sparingly: green (strong), amber (moderate), gray (none/expired)",
        ],
        font_size=16,
    )

    return prs


def main():
    deck = build_deck()
    deck.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()

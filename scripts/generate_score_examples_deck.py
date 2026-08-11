#!/usr/bin/env python3
"""Rebuild Score Examples.pptx with a clear, step-by-step lead scoring layout."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

WORKSPACE = Path(__file__).resolve().parents[1]
OUTPUT = WORKSPACE / "Score Examples.pptx"

CISCO_DARK = RGBColor(0x00, 0x52, 0x8B)
CISCO_BLUE = RGBColor(0x00, 0x7B, 0xC7)
CISCO_TEAL = RGBColor(0x00, 0x96, 0xA6)
CISCO_GREEN = RGBColor(0x6B, 0xB3, 0x45)
CISCO_AMBER = RGBColor(0xF5, 0xA6, 0x23)
CISCO_RED = RGBColor(0xCF, 0x3D, 0x3E)
CISCO_GRAY = RGBColor(0x58, 0x59, 0x5B)
CISCO_LIGHT = RGBColor(0xF2, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD0, 0xD7, 0xDE)

ENGAGEMENT_GRADES = [
    ("4", "76 – 100", "Hot — immediate follow-up"),
    ("3", "51 – 75", "Warm-high — prioritise"),
    ("2", "26 – 50", "Warm — nurture and monitor"),
    ("1", "0 – 25", "Cold — low-touch"),
]


def set_run(font, *, size: int, bold: bool = False, color: RGBColor | None = None):
    font.name = "Segoe UI"
    font.size = Pt(size)
    font.bold = bold
    if color:
        font.color.rgb = color


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CISCO_DARK
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_run(p.font, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.45), Inches(0.68), Inches(12.4), Inches(0.3))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_run(sp.font, size=13, color=WHITE)


def add_table(slide, left, top, width, height, headers, rows, col_widths=None, font_size=10):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = CISCO_DARK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            set_run(p.font, size=font_size, bold=True, color=WHITE)

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CISCO_LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c_idx >= 2 else PP_ALIGN.LEFT
                set_run(p.font, size=font_size, color=CISCO_GRAY)
    return table_shape


def add_outcome_banner(slide, left, top, width, text: str, accent: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_run(p.font, size=14, bold=True, color=WHITE)


def add_score_summary(
    slide,
    left,
    top,
    total: float,
    level: str,
    detail: str,
    banner_text: str,
    outcome_color: RGBColor,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.6), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = CISCO_LIGHT
    card.line.color.rgb = BORDER

    tf = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(3.2), Inches(1.7)).text_frame
    tf.clear()
    lines = [
        (f"Total engagement score: {total:g}", 16, True, CISCO_DARK),
        (f"Engagement level: {level}", 15, True, CISCO_BLUE),
        ("", 6, False, CISCO_GRAY),
        (detail, 13, False, CISCO_GRAY),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(4)
        set_run(p.font, size=size, bold=bold, color=color)

    add_outcome_banner(slide, left, top + Inches(2.15), Inches(3.6), banner_text, outcome_color)


def add_formula_box(slide, left, top, width, height):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = CISCO_BLUE

    tf = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3)).text_frame
    tf.clear()
    lines = [
        ("How each activity is scored", 16, True, CISCO_DARK),
        ("", 4, False, CISCO_GRAY),
        ("Points = Activity score % × Category weight %", 14, True, CISCO_BLUE),
        ("", 4, False, CISCO_GRAY),
        ("Example: 90% score × 15% weight = 13.5 points", 12, False, CISCO_GRAY),
        ("All qualifying activities in the look-back window are summed.", 12, False, CISCO_GRAY),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.space_after = Pt(4)
        set_run(p.font, size=size, bold=bold, color=color)


def add_timeline_step_table(slide, steps: list[dict]):
    """Add a step-by-step scoring table below the slide title bar."""
    headers = ["Step", "Activity & rule", "Score %", "Weight", "Points", "Running total"]
    rows = []
    for step in steps:
        rows.append(
            [
                step["step"],
                step["activity"],
                step["score_pct"],
                step["weight"],
                step["points"],
                step["running_total"],
            ]
        )

    add_table(
        slide,
        Inches(0.45),
        Inches(1.25),
        Inches(9.0),
        Inches(min(5.8, 0.45 * (len(rows) + 1))),
        headers,
        rows,
        col_widths=[
            Inches(0.55),
            Inches(3.8),
            Inches(0.85),
            Inches(0.85),
            Inches(0.85),
            Inches(1.1),
        ],
    )

    return rows


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CISCO_DARK
    bg.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Lead Scoring Examples"
    set_run(tp.font, size=40, bold=True, color=WHITE)

    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.0))
    stf = subtitle.text_frame
    stf.clear()
    for i, line in enumerate(
        [
            "Detailed walkthroughs showing how engagement points are calculated",
            "Amanda Chenery  •  Marketing Operations",
        ]
    ):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        set_run(p.font, size=18 if i == 0 else 14, color=WHITE)

    # Slide 2 — How scoring works
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "How Engagement Scoring Works", "Every activity contributes points; the total determines the engagement level")

    add_formula_box(slide, Inches(0.45), Inches(1.3), Inches(5.5), Inches(2.2))

    add_table(
        slide,
        Inches(6.3),
        Inches(1.3),
        Inches(6.5),
        Inches(2.2),
        ["Level", "Score range", "Meaning"],
        ENGAGEMENT_GRADES,
        col_widths=[Inches(0.7), Inches(1.3), Inches(4.0)],
    )

    add_bullets = slide.shapes.add_textbox(Inches(0.45), Inches(3.8), Inches(12.4), Inches(3.0)).text_frame
    add_bullets.clear()
    bullet_lines = [
        "Key principles used in every example on the following slides:",
        "• Only activities within the recency window (7, 14, or 30 days) contribute points.",
        "• Higher-intent activities carry a larger category weight (e.g. demo = 25%, email click = 5%).",
        "• Multiple activities stack — the running total is the sum of all qualifying points.",
        "• Leads are routed to VDC when engagement is high enough AND profile fit is A, B, or C.",
    ]
    for i, line in enumerate(bullet_lines):
        p = add_bullets.paragraphs[0] if i == 0 else add_bullets.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
        set_run(p.font, size=14 if i == 0 else 13, bold=i == 0, color=CISCO_DARK if i == 0 else CISCO_GRAY)

    # Slide 3 — Example 1: First touch content syndication
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 1: First-Touch Content Syndication",
        "Contact enters Eloqua with only syndicated content activity — low score, no VDC routing",
    )

    steps = [
        {
            "step": "1",
            "activity": "Form submit from Integrate (non-hand-raiser) — at least 1 time in last 7 days",
            "score_pct": "90%",
            "weight": "10%",
            "points": "9.0",
            "running_total": "9.0",
        },
        {
            "step": "2",
            "activity": "Content syndication activity created — in last 7 days",
            "score_pct": "90%",
            "weight": "15%",
            "points": "13.5",
            "running_total": "22.5",
        },
    ]
    add_timeline_step_table(slide, steps)

    add_score_summary(
        slide,
        Inches(9.7),
        Inches(1.25),
        22.5,
        "1",
        "Profile must also be A/B/C for any VDC routing.",
        "Low scored — no further action",
        CISCO_RED,
    )

    note = slide.shapes.add_textbox(Inches(0.45), Inches(3.2), Inches(9.0), Inches(1.2)).text_frame
    note.clear()
    for i, line in enumerate(
        [
            "Calculation check: 9.0 + 13.5 = 22.5",
            "Engagement level 1 (0–25) = cold. This lead stays in marketing nurture.",
        ]
    ):
        p = note.paragraphs[0] if i == 0 else note.add_paragraph()
        p.text = line
        set_run(p.font, size=12, bold=i == 0, color=CISCO_DARK if i == 0 else CISCO_GRAY)

    # Slide 4 — Example 2: Multi-touch
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 2: Multi-Touch with Content Syndication Last Touch",
        "Prior digital activity stacks; last-touch syndication pushes the score above the VDC threshold",
    )

    multi_steps = [
        ("1", "Email click-through — in last 14 days", "80%", "5%", 4.0),
        ("2", "Website visit — in last 14 days", "80%", "10%", 8.0),
        ("3", "Landing page visit — in last 14 days", "80%", "15%", 12.0),
        ("4", "Video watched — in last 7 days", "100%", "15%", 15.0),
        ("5", "Webinar watched — in last 7 days", "90%", "10%", 9.0),
        ("6", "Content syndication lead — in last 7 days", "90%", "15%", 13.5),
    ]
    running = 0.0
    rows = []
    for step_num, activity, score_pct, weight, points in multi_steps:
        running += points
        rows.append(
            {
                "step": step_num,
                "activity": activity,
                "score_pct": score_pct,
                "weight": weight,
                "points": f"{points:g}",
                "running_total": f"{running:g}",
            }
        )

    add_timeline_step_table(slide, rows)

    add_score_summary(
        slide,
        Inches(9.7),
        Inches(1.25),
        61.5,
        "2",
        "Requires profile grade A, B, or C.",
        "High scored — passed to VDC",
        CISCO_GREEN,
    )

    # Slide 5 — Event registration step-by-step
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 3: Event Registration — Step by Step",
        "Each registration adds points; score accumulates as the contact registers for more sessions",
    )

    reg_steps = [
        (
            "1",
            "8 May 2024 — Registers for Cisco Live (CLUS)",
            "Registered at least 1 time in last 7 days",
            "90%",
            "10%",
            "9",
            "9",
        ),
        (
            "2",
            "9 May 2024 — Registers for Keynote session",
            "Registered more than 1 time in last 7 days",
            "100%",
            "10%",
            "10",
            "10",
        ),
        (
            "3",
            "16 May 2024 — Registers for breakout session",
            "Registered more than 1 time in last 7 days",
            "100%",
            "10%",
            "10",
            "10",
        ),
    ]

    add_table(
        slide,
        Inches(0.45),
        Inches(1.25),
        Inches(12.4),
        Inches(2.5),
        ["Step", "Date & event", "Scoring rule", "Score %", "Weight", "Points", "Running total"],
        reg_steps,
        col_widths=[
            Inches(0.5),
            Inches(2.8),
            Inches(3.2),
            Inches(0.75),
            Inches(0.75),
            Inches(0.75),
            Inches(1.0),
        ],
        font_size=9,
    )

    add_outcome_banner(
        slide,
        Inches(0.45),
        Inches(4.0),
        Inches(12.4),
        "After step 3: Total = 10 points  •  Engagement level 1  •  No further action (registrations alone)",
        CISCO_AMBER,
    )

    logic = slide.shapes.add_textbox(Inches(0.45), Inches(4.9), Inches(12.4), Inches(2.0)).text_frame
    logic.clear()
    for i, line in enumerate(
        [
            "Logic → Eloqua identifies each registration transaction and matches it to the Event Registration category.",
            "Calculation → Points = score % × 10% weight. Frequency rules (at least 1 vs. more than 1) determine the score %.",
            "Action → Score accumulates but remains below the high-score threshold until attendance or additional activity occurs.",
        ]
    ):
        p = logic.paragraphs[0] if i == 0 else logic.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        set_run(p.font, size=13, color=CISCO_GRAY)

    # Slide 6 — Event attendance step-by-step
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 4: Event Attendance — Step by Step",
        "Attendance carries a higher category weight (30%) and can trigger VDC routing",
    )

    att_steps = [
        (
            "1",
            "2 Jun 2024 — Attends Cisco Live (CLUS)",
            "Attended at least 1 time in last 7 days",
            "90%",
            "30%",
            "27",
            "27",
        ),
        (
            "2",
            "3 Jun 2024 — Attends Keynote session",
            "Attended more than 1 time in last 7 days",
            "100%",
            "30%",
            "30",
            "30",
        ),
        (
            "3",
            "4 Jun 2024 — Attends breakout session",
            "Attended more than 1 time in last 7 days",
            "100%",
            "30%",
            "30",
            "30",
        ),
    ]

    add_table(
        slide,
        Inches(0.45),
        Inches(1.25),
        Inches(12.4),
        Inches(2.5),
        ["Step", "Date & event", "Scoring rule", "Score %", "Weight", "Points", "Running total"],
        att_steps,
        col_widths=[
            Inches(0.5),
            Inches(2.8),
            Inches(3.2),
            Inches(0.75),
            Inches(0.75),
            Inches(0.75),
            Inches(1.0),
        ],
        font_size=9,
    )

    add_outcome_banner(
        slide,
        Inches(0.45),
        Inches(4.0),
        Inches(12.4),
        "After step 3: Total = 30 points  •  Engagement level 2  •  High-score lead if profile is A",
        CISCO_GREEN,
    )

    # Slide 7 — Decay on older registrations
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 5: Why Older Registrations Still Count (with Decay)",
        "Registrations from the last 30 days are included, but points are lower outside the 7-day window",
    )

    decay_note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(12.4), Inches(0.9))
    decay_note.fill.solid()
    decay_note.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
    decay_note.line.color.rgb = CISCO_AMBER
    dp = slide.shapes.add_textbox(Inches(0.65), Inches(1.4), Inches(12.0), Inches(0.7)).text_frame.paragraphs[0]
    dp.text = (
        "When a contact attends an event, Eloqua also re-scores earlier registrations from the last 30 days "
        "— but at reduced score percentages due to recency decay."
    )
    set_run(dp.font, size=13, color=CISCO_GRAY)

    decay_rows = [
        ("Registration", "8 May — CLUS", "At least 1 time, 30 days", "50%", "10%", "5", "5"),
        ("Registration", "9 May — Keynote", "More than 1 time, 30 days", "60%", "10%", "6", "6"),
        ("Registration", "16 May — Breakout", "More than 1 time, 30 days", "60%", "10%", "6", "6"),
        ("Attendance", "2 Jun — CLUS", "At least 1 time, 7 days", "90%", "30%", "27", "33"),
        ("Attendance", "3 Jun — Keynote", "More than 1 time, 7 days", "100%", "30%", "30", "36"),
        ("Attendance", "4 Jun — Breakout", "More than 1 time, 7 days", "100%", "30%", "30", "36"),
    ]

    add_table(
        slide,
        Inches(0.45),
        Inches(2.35),
        Inches(12.4),
        Inches(3.2),
        ["Type", "Activity", "Scoring rule", "Score %", "Weight", "Points", "Running total"],
        decay_rows,
        col_widths=[
            Inches(1.1),
            Inches(2.2),
            Inches(2.8),
            Inches(0.75),
            Inches(0.75),
            Inches(0.75),
            Inches(1.0),
        ],
        font_size=9,
    )

    foot = slide.shapes.add_textbox(Inches(0.45), Inches(5.75), Inches(12.4), Inches(0.8)).text_frame.paragraphs[0]
    foot.text = (
        "Older registrations still contribute at reduced 30-day decay rates. Fresh attendance at 7-day "
        "rates pushes the total to 36. Engagement level 3 → passed to VDC if profile is A."
    )
    set_run(foot.font, size=11, color=CISCO_GRAY)

    # Slide 8 — Email click driving registration
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 6: Digital Touchpoints Before Event Registration",
        "Email and web activity that preceded registration is also included in the final score",
    )

    digital_note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(1.25), Inches(12.4), Inches(0.9))
    digital_note.fill.solid()
    digital_note.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
    digital_note.line.color.rgb = CISCO_AMBER
    dnp = slide.shapes.add_textbox(Inches(0.65), Inches(1.4), Inches(12.0), Inches(0.7)).text_frame.paragraphs[0]
    dnp.text = (
        "If the contact clicked an email to reach Cisco.com before registering, those transactions "
        "within the last 30 days are included alongside registrations and attendance."
    )
    set_run(dnp.font, size=13, color=CISCO_GRAY)

    digital_rows = [
        ("Digital", "8 May — Email click-through", "At least 1 time, 30 days", "50%", "5%", "2.5", "2.5"),
        ("Digital", "8 May — Landing page visit", "At least 1 time, 30 days", "50%", "15%", "7.5", "10.0"),
        ("Digital", "8 May — Website visit", "At least 1 time, 30 days", "50%", "10%", "5.0", "15.0"),
        ("Registration", "8 May — CLUS", "At least 1 time, 30 days", "50%", "10%", "5.0", "20.0"),
        ("Registration", "9 May — Keynote", "More than 1 time, 30 days", "60%", "10%", "6.0", "21.0"),
        ("Registration", "16 May — Breakout", "More than 1 time, 30 days", "60%", "10%", "6.0", "21.0"),
    ]

    add_table(
        slide,
        Inches(0.45),
        Inches(2.35),
        Inches(12.4),
        Inches(3.0),
        ["Type", "Activity", "Scoring rule", "Score %", "Weight", "Points", "Running total"],
        digital_rows,
        col_widths=[
            Inches(1.1),
            Inches(2.2),
            Inches(2.8),
            Inches(0.75),
            Inches(0.75),
            Inches(0.75),
            Inches(1.0),
        ],
        font_size=9,
    )

    # Slide 9 — Full view
    slide = prs.slides.add_slide(blank)
    add_title_bar(
        slide,
        "Example 7: Full Journey — All Touchpoints Combined",
        "Complete timeline from first email click through event attendance",
    )

    full_rows = [
        ("Digital", "8 May — Email click", "At least 1, 30 days", "50%", "5%", "2.5", "2.5"),
        ("Digital", "8 May — Landing page", "At least 1, 30 days", "50%", "15%", "7.5", "10.0"),
        ("Digital", "8 May — Website visit", "At least 1, 30 days", "50%", "10%", "5.0", "15.0"),
        ("Registration", "8 May — CLUS", "At least 1, 30 days", "50%", "10%", "5.0", "20.0"),
        ("Registration", "9 May — Keynote", "More than 1, 30 days", "60%", "10%", "6.0", "21.0"),
        ("Registration", "16 May — Breakout", "More than 1, 30 days", "60%", "10%", "6.0", "21.0"),
        ("Attendance", "2 Jun — CLUS", "At least 1, 7 days", "90%", "30%", "27.0", "48.0"),
        ("Attendance", "3 Jun — Keynote", "More than 1, 7 days", "100%", "30%", "30.0", "51.0"),
        ("Attendance", "4 Jun — Breakout", "More than 1, 7 days", "100%", "30%", "30.0", "51.0"),
    ]

    add_table(
        slide,
        Inches(0.45),
        Inches(1.25),
        Inches(12.4),
        Inches(4.5),
        ["Type", "Activity", "Scoring rule", "Score %", "Weight", "Points", "Running total"],
        full_rows,
        col_widths=[
            Inches(1.0),
            Inches(2.0),
            Inches(2.6),
            Inches(0.7),
            Inches(0.7),
            Inches(0.7),
            Inches(0.95),
        ],
        font_size=8,
    )

    add_outcome_banner(
        slide,
        Inches(0.45),
        Inches(5.95),
        Inches(12.4),
        "Final total = 51  •  Engagement level 2  •  Passed to VDC if profile is A, B, or C",
        CISCO_GREEN,
    )

    summary = slide.shapes.add_textbox(Inches(0.45), Inches(6.75), Inches(12.4), Inches(0.5)).text_frame.paragraphs[0]
    summary.text = (
        "This matches the original full-view example: digital + registration (30-day decay) + attendance (7-day) "
        "all stack to produce the final engagement score."
    )
    set_run(summary.font, size=11, color=CISCO_GRAY)

    # Slide 10 — Summary / routing rules
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, "Summary: When Does a Lead Route to VDC?", "Engagement score and profile grade must both qualify")

    add_table(
        slide,
        Inches(0.8),
        Inches(1.5),
        Inches(5.5),
        Inches(2.0),
        ["Example", "Engagement total", "Level", "Outcome"],
        [
            ["First-touch syndication", "22.5", "1", "No action"],
            ["Multi-touch syndication", "61.5", "2", "Passed to VDC*"],
            ["Event registration only", "10", "1", "No action"],
            ["Event attendance", "30", "2", "High score if profile A*"],
            ["Full journey", "51", "2", "Passed to VDC*"],
        ],
        col_widths=[Inches(2.2), Inches(1.2), Inches(0.7), Inches(1.4)],
    )

    routing = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.0)).text_frame
    routing.clear()
    routing_lines = [
        ("Routing rules", 18, True, CISCO_DARK),
        ("", 6, False, CISCO_GRAY),
        ("* VDC routing also requires:", 14, True, CISCO_BLUE),
        ("• Profile grade A, B, or C", 13, False, CISCO_GRAY),
        ("• Valid CCID present", 13, False, CISCO_GRAY),
        ("• Not AU Public Sector (unless exempt)", 13, False, CISCO_GRAY),
        ("• Passes Program 389 validation", 13, False, CISCO_GRAY),
        ("", 6, False, CISCO_GRAY),
        ("Remember:", 14, True, CISCO_DARK),
        ("Points = Score % × Weight %", 13, False, CISCO_GRAY),
        ("All qualifying activities in the look-back window sum to the total.", 13, False, CISCO_GRAY),
        ("Recency decay reduces points for older activity (30-day < 14-day < 7-day).", 13, False, CISCO_GRAY),
    ]
    for i, (text, size, bold, color) in enumerate(routing_lines):
        p = routing.paragraphs[0] if i == 0 else routing.add_paragraph()
        p.text = text
        p.space_after = Pt(5)
        set_run(p.font, size=size, bold=bold, color=color)

    return prs


def main():
    deck = build_deck()
    deck.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
